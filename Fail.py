import telebot
from telebot import types
from PIL import Image, ImageDraw, ImageFont
from docx import Document
import openpyxl
from pptx import Presentation
from io import BytesIO
import os, tempfile, threading
import qrcode
import tempfile
import threading
from io import BytesIO

# ======== توکن ربات ========
TOKEN = "5902797816:AAFYD_WH1k-cvB2T2NvMW3t8Zitv7tU2Bhc"
bot = telebot.TeleBot(TOKEN)

# ---------- ذخیره داده‌های کاربران ----------
user_photos = {}
user_text_mode = set()  # کاربرانی که در حالت متن هستند

# ---------- منوی اصلی ----------
def main_menu():
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    markup.row('📸 عکس به PDF', '📄 Word به PDF')
    markup.row('📊 Excel به PDF', '🎞 PowerPoint به PDF')
    markup.row('📝 متن به PDF', '🎭 استیکر به عکس')
    markup.row('🟢 ساخت QR کد')
    markup.row('📌 راهنما')
    return markup

# ---------- شروع ----------
@bot.message_handler(commands=['start'])
def send_welcome(message):
    bot.send_message(message.chat.id,
                     "سلام! 🤖\nمن ربات همه‌کاره تبدیل فایل هستم.\nیکی از گزینه‌های منو رو انتخاب کن.",
                     reply_markup=main_menu())

# ---------- راهنما ----------
@bot.message_handler(func=lambda m: m.text == '📌 راهنما')
def send_help(message):
    bot.send_message(message.chat.id,
                     "📖 راهنما:\n"
                     "- عکس‌ها رو به PDF تبدیل کن 📸\n"
                     "- فایل Word، Excel و PowerPoint رو به PDF کن 📄📊🎞\n"
                     "- متن ساده رو PDF بگیر 📝\n"
                     "- استیکر رو عکس کن 🎭\n\n"
                     "✅ فایل‌ها بعد از ۳۰ ثانیه پاک می‌شن.",
                     reply_markup=main_menu())

# ---------- عکس به PDF ----------
@bot.message_handler(func=lambda m: m.text == '📸 عکس به PDF')
def start_photo_pdf(message):
    user_photos[message.chat.id] = []
    bot.send_message(message.chat.id, "📸 عکس‌هاتو یکی یکی بفرست و وقتی تموم شد بگو «تموم شد».")

@bot.message_handler(content_types=['photo'])
def handle_photos(message):
    if message.chat.id in user_photos:
        file_info = bot.get_file(message.photo[-1].file_id)
        downloaded = bot.download_file(file_info.file_path)
        img = Image.open(BytesIO(downloaded))
        user_photos[message.chat.id].append(img)
        bot.send_message(message.chat.id, f"✅ عکس ذخیره شد ({len(user_photos[message.chat.id])})")

@bot.message_handler(func=lambda m: m.text.lower() == 'تموم شد')
def finish_photos(message):
    chat_id = message.chat.id
    if chat_id not in user_photos or not user_photos[chat_id]:
        bot.send_message(chat_id, "❌ عکسی ذخیره نشده!")
        return
    pdf_path = tempfile.mktemp(".pdf")
    user_photos[chat_id][0].save(pdf_path, save_all=True, append_images=user_photos[chat_id][1:])
    del user_photos[chat_id]
    bot.send_document(message.chat.id, open(pdf_path, "rb"), caption="📄 PDF آماده شد ✅")
    threading.Timer(30, lambda: os.remove(pdf_path)).start()

# ---------- Word به PDF ----------
@bot.message_handler(func=lambda m: m.text == '📄 Word به PDF')
def start_word(message):
    bot.send_message(message.chat.id, "📄 فایل Word رو بفرست.")

@bot.message_handler(content_types=['document'])
def handle_docs(message):
    if message.document.file_name.endswith(".docx"):
        file_info = bot.get_file(message.document.file_id)
        downloaded = bot.download_file(file_info.file_path)
        doc = Document(BytesIO(downloaded))
        lines = [p.text for p in doc.paragraphs]
        img = Image.new("RGB", (800, 1000), "white")
        d = ImageDraw.Draw(img)
        font = ImageFont.load_default()
        y = 10
        pages = []
        for line in lines:
            d.text((10, y), line, font=font, fill="black")
            y += 20
            if y > 950:
                pages.append(img)
                img = Image.new("RGB", (800, 1000), "white")
                d = ImageDraw.Draw(img)
                y = 10
        pages.append(img)
        pdf_path = tempfile.mktemp(".pdf")
        pages[0].save(pdf_path, save_all=True, append_images=pages[1:])
        bot.send_document(message.chat.id, open(pdf_path, "rb"), caption="📄 Word تبدیل شد ✅")
        threading.Timer(30, lambda: os.remove(pdf_path)).start()

# ---------- Excel به PDF ----------
@bot.message_handler(func=lambda m: m.text == '📊 Excel به PDF')
def start_excel(message):
    bot.send_message(message.chat.id, "📊 فایل Excel رو بفرست.")

@bot.message_handler(func=lambda m: m.document and m.document.file_name.endswith(".xlsx"))
def handle_excel(message):
    file_info = bot.get_file(message.document.file_id)
    downloaded = bot.download_file(file_info.file_path)
    wb = openpyxl.load_workbook(BytesIO(downloaded))
    sheet = wb.active
    lines = [[str(cell.value) for cell in row] for row in sheet.iter_rows()]
    img = Image.new("RGB", (1000, 1400), "white")
    d = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    y = 10
    for row in lines:
        d.text((10, y), " | ".join(row), font=font, fill="black")
        y += 20
    pdf_path = tempfile.mktemp(".pdf")
    img.save(pdf_path)
    bot.send_document(message.chat.id, open(pdf_path, "rb"), caption="📊 Excel تبدیل شد ✅")
    threading.Timer(30, lambda: os.remove(pdf_path)).start()

# ---------- PowerPoint به PDF ----------
@bot.message_handler(func=lambda m: m.text == '🎞 PowerPoint به PDF')
def start_ppt(message):
    bot.send_message(message.chat.id, "🎞 فایل PowerPoint رو بفرست.")

@bot.message_handler(func=lambda m: m.document and m.document.file_name.endswith(".pptx"))
def handle_ppt(message):
    file_info = bot.get_file(message.document.file_id)
    downloaded = bot.download_file(file_info.file_path)
    prs = Presentation(BytesIO(downloaded))
    slides = []
    for slide in prs.slides:
        img = Image.new("RGB", (800, 600), "white")
        d = ImageDraw.Draw(img)
        y = 10
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                d.text((10, y), shape.text, fill="black")
                y += 20
        slides.append(img)
    pdf_path = tempfile.mktemp(".pdf")
    slides[0].save(pdf_path, save_all=True, append_images=slides[1:])
    bot.send_document(message.chat.id, open(pdf_path, "rb"), caption="🎞 پاورپوینت تبدیل شد ✅")
    threading.Timer(30, lambda: os.remove(pdf_path)).start()

# ---------- متن به PDF ----------
@bot.message_handler(func=lambda m: m.text == '📝 متن به PDF')
def start_text_pdf(message):
    user_text_mode.add(message.chat.id)
    bot.send_message(message.chat.id, "📝 متن خودت رو اینجا بفرست. بعد از ارسال، PDF ساخته می‌شه.")

@bot.message_handler(func=lambda m: m.chat.id in user_text_mode)
def handle_text_pdf(message):
    text = message.text
    img = Image.new("RGB", (800, 1000), "white")
    d = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    y = 10
    for line in text.split("\n"):
        d.text((10, y), line, font=font, fill="black")
        y += 20
    pdf_path = tempfile.mktemp(".pdf")
    img.save(pdf_path)
    bot.send_document(message.chat.id, open(pdf_path, "rb"), caption="📝 PDF آماده شد ✅")
    threading.Timer(30, lambda: os.remove(pdf_path)).start()
    user_text_mode.remove(message.chat.id)

# ---------- استیکر به عکس ----------
@bot.message_handler(func=lambda m: m.text == '🎭 استیکر به عکس')
def start_sticker(message):
    bot.send_message(message.chat.id, "🎭 لطفاً استیکر ثابت خود را ارسال کنید:")

@bot.message_handler(content_types=['sticker'])
def handle_sticker(message):
    try:
        file_info = bot.get_file(message.sticker.file_id)
        downloaded = bot.download_file(file_info.file_path)
        img = Image.open(BytesIO(downloaded))
        if getattr(img, "is_animated", False):
            bot.send_message(message.chat.id, "❌ این استیکر متحرک است و قابل تبدیل نیست.")
            return
        path = tempfile.mktemp(".png")
        img.save(path)
        bot.send_document(message.chat.id, open(path, "rb"), caption="🎭 استیکر تبدیل شد ✅")
        threading.Timer(30, lambda: os.remove(path)).start()
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ خطا در تبدیل استیکر: {str(e)}")

# ---------- QR کد ----------
user_qr_mode = set()

@bot.message_handler(func=lambda m: m.text == '🟢 QR کد بساز')
def start_qr(message):
    # پاک کردن حالت‌های دیگر کاربر
    user_photos.pop(message.chat.id, None)
    user_text_mode.discard(message.chat.id)
    
    user_qr_mode.add(message.chat.id)
    bot.send_message(message.chat.id, "📌 متن یا لینک خودت رو بفرست تا QR کد بسازم:")

@bot.message_handler(func=lambda m: m.chat.id in user_qr_mode)
def generate_qr(message):
    # فقط پیام متنی پردازش بشه
    if message.content_type != 'text':
        bot.send_message(message.chat.id, "❌ لطفا فقط متن یا لینک وارد کنید.")
        return

    text = message.text
    img = qrcode.make(text)
    path = tempfile.mktemp(".png")
    img.save(path)
    
    bot.send_document(message.chat.id, open(path, "rb"), caption="✅ QR کد آماده شد")
    threading.Timer(30, lambda: os.remove(path)).start()
    
    user_qr_mode.remove(message.chat.id)

# ---------- اجرای ربات ----------
bot.infinity_polling()
